"""
PowerPoint Report Generation Pipeline Step
GroundTruth Hackathon | Automated Insight Engine

This module handles:
- Creating professional PPTX presentations using python-pptx
- Including cover slide, KPI slides, charts, and insights
"""

import os
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Brand colors (RGB) - using RGBColor (capital letters)
BRAND_PRIMARY = RGBColor(37, 99, 235)     # #2563EB
BRAND_SECONDARY = RGBColor(16, 185, 129)  # #10B981
BRAND_DARK = RGBColor(30, 41, 59)         # #1E293B
BRAND_LIGHT = RGBColor(241, 245, 249)     # #F1F5F9
WHITE = RGBColor(255, 255, 255)


def create_title_slide(prs: Presentation, title: str, report_id: str) -> None:
    """Create the title/cover slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BRAND_PRIMARY
    background.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(9), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "Automated Insight Engine Report"
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add date
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5),
        Inches(9), Inches(0.5)
    )
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = datetime.now().strftime("%B %d, %Y")
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = WHITE
    date_para.alignment = PP_ALIGN.CENTER
    
    # Add report ID
    id_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5),
        Inches(9), Inches(0.3)
    )
    id_frame = id_box.text_frame
    id_para = id_frame.paragraphs[0]
    id_para.text = f"Report ID: {report_id[:8]}..."
    id_para.font.size = Pt(10)
    id_para.font.color.rgb = WHITE
    id_para.alignment = PP_ALIGN.CENTER


def create_section_slide(prs: Presentation, section_title: str) -> None:
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(3),
        Inches(10), Inches(1.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_PRIMARY
    accent.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.25),
        Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = section_title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER


def create_executive_summary_slide(
    prs: Presentation,
    insights: Dict[str, Any]
) -> None:
    """Create executive summary slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Executive Summary"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY
    
    # Add summary text
    summary = insights.get('executive_summary', 'No summary available.')
    
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content_para = content_frame.paragraphs[0]
    content_para.text = summary
    content_para.font.size = Pt(18)
    content_para.font.color.rgb = BRAND_DARK
    content_para.line_spacing = 1.5


def create_kpi_slide(prs: Presentation, kpis: Dict[str, Any]) -> None:
    """Create KPI overview slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Key Performance Indicators"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY
    
    basic_metrics = kpis.get('basic_metrics', {})
    
    # Define KPIs to display
    kpi_items = [
        ('Total Impressions', 'total_impressions', '{:,.0f}'),
        ('Total Clicks', 'total_clicks', '{:,.0f}'),
        ('Overall CTR', 'overall_ctr', '{:.2f}%'),
        ('Total Spend', 'total_spend', '${:,.2f}'),
        ('Total Conversions', 'total_conversions', '{:,.0f}'),
        ('Conversion Rate', 'overall_conversion_rate', '{:.2f}%'),
    ]
    
    # Create KPI boxes in a grid
    start_x = 0.5
    start_y = 1.5
    box_width = 2.8
    box_height = 1.8
    gap = 0.3
    
    for i, (label, key, fmt) in enumerate(kpi_items):
        if key not in basic_metrics:
            continue
            
        row = i // 3
        col = i % 3
        
        x = start_x + col * (box_width + gap)
        y = start_y + row * (box_height + gap)
        
        # Add KPI box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = BRAND_LIGHT
        box.line.color.rgb = BRAND_PRIMARY
        
        # Add value
        value = basic_metrics[key]
        try:
            formatted_value = fmt.format(value)
        except:
            formatted_value = str(value)
        
        value_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.3),
            Inches(box_width), Inches(0.8)
        )
        value_frame = value_box.text_frame
        value_para = value_frame.paragraphs[0]
        value_para.text = formatted_value
        value_para.font.size = Pt(28)
        value_para.font.bold = True
        value_para.font.color.rgb = BRAND_PRIMARY
        value_para.alignment = PP_ALIGN.CENTER
        
        # Add label
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + 1.1),
            Inches(box_width), Inches(0.5)
        )
        label_frame = label_box.text_frame
        label_para = label_frame.paragraphs[0]
        label_para.text = label
        label_para.font.size = Pt(12)
        label_para.font.color.rgb = BRAND_DARK
        label_para.alignment = PP_ALIGN.CENTER


def create_highlights_slide(
    prs: Presentation,
    insights: Dict[str, Any]
) -> None:
    """Create key highlights slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Key Highlights"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY
    
    highlights = insights.get('key_highlights', [])
    
    # Add highlights
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, highlight in enumerate(highlights[:5]):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        
        para.text = f"✓ {highlight}"
        para.font.size = Pt(16)
        para.font.color.rgb = BRAND_DARK
        para.space_after = Pt(12)
        para.level = 0


def create_issues_slide(
    prs: Presentation,
    insights: Dict[str, Any]
) -> None:
    """Create performance issues slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Performance Issues"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY
    
    issues = insights.get('performance_issues', [])
    
    # Add issues
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, issue in enumerate(issues[:3]):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        
        para.text = f"⚠ {issue}"
        para.font.size = Pt(16)
        para.font.color.rgb = BRAND_DARK
        para.space_after = Pt(12)


def create_recommendations_slide(
    prs: Presentation,
    insights: Dict[str, Any]
) -> None:
    """Create recommendations slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Recommendations"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY
    
    recommendations = insights.get('recommendations', [])
    
    # Add recommendations
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, rec in enumerate(recommendations[:3]):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        
        para.text = f"{i + 1}. {rec}"
        para.font.size = Pt(16)
        para.font.color.rgb = BRAND_DARK
        para.space_after = Pt(16)


def create_chart_slide(
    prs: Presentation,
    chart_path: str,
    chart_title: str
) -> None:
    """Create a slide with a chart image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = chart_title
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_PRIMARY
    
    # Add chart image
    if os.path.exists(chart_path):
        slide.shapes.add_picture(
            chart_path,
            Inches(0.5), Inches(1.2),
            width=Inches(9), height=Inches(5.5)
        )


def create_thank_you_slide(prs: Presentation) -> None:
    """Create thank you/closing slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BRAND_PRIMARY
    background.line.fill.background()
    
    # Add thank you text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3),
        Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Thank You"
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5),
        Inches(9), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "Generated by Automated Insight Engine"
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER


def generate_ppt_report(
    report_id: str,
    report_title: str,
    kpis: Dict[str, Any],
    charts: Dict[str, Any],
    insights: Dict[str, Any],
    output_dir: str
) -> Dict[str, Any]:
    """
    Main PowerPoint report generation function.
    
    This function creates a professional PPTX presentation including:
    1. Title slide
    2. Executive summary
    3. KPI overview
    4. Key highlights
    5. Performance issues
    6. Recommendations
    7. Chart slides
    8. Thank you slide
    
    Args:
        report_id: Unique report identifier
        report_title: Title for the report
        kpis: Dictionary of computed KPIs
        charts: Dictionary containing chart paths
        insights: Dictionary containing LLM-generated insights
        output_dir: Directory to save the report
        
    Returns:
        Dictionary with file path and status
    """
    logger.info(f"Generating PowerPoint report: {report_id}")
    
    # Create output directory
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Define output file path
    file_path = output_path / f"{report_id}.pptx"
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Build slides
    try:
        # Title slide
        create_title_slide(prs, report_title, report_id)
        
        # Executive summary
        create_executive_summary_slide(prs, insights)
        
        # KPI overview
        create_kpi_slide(prs, kpis)
        
        # Key highlights
        create_highlights_slide(prs, insights)
        
        # Performance issues
        create_issues_slide(prs, insights)
        
        # Recommendations
        create_recommendations_slide(prs, insights)
        
        # Charts section
        chart_paths = charts.get('chart_paths', [])
        if chart_paths:
            create_section_slide(prs, "Visual Analytics")
            
            for chart_path in chart_paths:
                if os.path.exists(chart_path):
                    chart_name = Path(chart_path).stem.replace('_', ' ').title()
                    create_chart_slide(prs, chart_path, chart_name)
        
        # Thank you slide
        create_thank_you_slide(prs)
        
        # Save presentation
        prs.save(str(file_path))
        
        logger.info(f"PowerPoint report generated: {file_path}")
        
        return {
            'status': 'success',
            'file_path': str(file_path),
            'format': 'pptx',
            'slides': len(prs.slides)
        }
        
    except Exception as e:
        logger.error(f"Failed to generate PowerPoint: {str(e)}")
        raise

